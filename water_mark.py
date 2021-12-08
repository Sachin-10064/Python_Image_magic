from pptx import Presentation
from pptx.util import Inches
from wand.image import Image

image_path = ["image1.jpg", "image2.jpg", "image3.jpg", "image4.jpg", "image5.jpg"]

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

for img in image_path:
    logo = Image(filename="nike_black.png")
    image = Image(filename=img)
    logo.resize(int(image.width / 2 - 30), int(image.width / 6))
    image.composite_channel("all_channels",
                            logo,
                            "dissolve",
                            int(10),
                            int(10)
                            )
    # img_name = img.split(".")
    # img_path = img_name[0]+"1.png"
    # print(img_path)
    image.save(filename=img)

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Sample Title1'

    tf = body_shape.text_frame
    tf.text = 'Sample Subtitle1'

    left = Inches(1)
    top = Inches(2.4)
    height = Inches(4.5)
    pic = slide.shapes.add_picture(img, left, top, height=height)

prs.save('test.pptx')